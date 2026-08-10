"""Pedalshield investor deck.

Facts in this deck were verified against the live service and repo, not the
stale committed config (deploy/pedalshield-backend.service still carries the
old 793 zat/km rate; /treasury/info reports 10,596).
"""
import sys
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

BG    = "0B1512"
PANEL = "132420"
EDGE  = "1E322B"
MINT  = "2BD99F"
PURP  = "A78BFA"
INK   = "EAF6F0"
MUTE  = "93A8A0"
DIM   = "6E8079"

H = "Arial"
B = "Calibri"

M = 0.7                 # left margin, inches
SW, SH = 13.333, 7.5
W = SW - M * 2

prs = Presentation()
prs.slide_width = Inches(SW)
prs.slide_height = Inches(SH)
BLANK = prs.slide_layouts[6]


def rgb(h):
    return RGBColor.from_string(h)


def new_slide():
    s = prs.slides.add_slide(BLANK)
    fill = s.background.fill
    fill.solid()
    fill.fore_color.rgb = rgb(BG)
    return s


def text(slide, txt, x, y, w, h, *, size=14, color=MUTE, font=B, bold=False,
         italic=False, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
         spacing=1.0, space_after=0):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    tf.vertical_anchor = anchor
    for i, line in enumerate(str(txt).split("\n")):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.line_spacing = spacing
        if space_after:
            p.space_after = Pt(space_after)
        r = p.add_run()
        r.text = line
        r.font.size = Pt(size)
        r.font.bold = bold
        r.font.italic = italic
        r.font.name = font
        r.font.color.rgb = rgb(color)
    return box


def panel(slide, x, y, w, h, *, fill=PANEL, line=EDGE, line_w=1.0, radius=0.06):
    sh = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                Inches(x), Inches(y), Inches(w), Inches(h))
    sh.adjustments[0] = radius
    sh.fill.solid()
    sh.fill.fore_color.rgb = rgb(fill)
    sh.line.color.rgb = rgb(line)
    sh.line.width = Pt(line_w)
    sh.shadow.inherit = False
    if sh.has_text_frame:
        sh.text_frame.text = ""
    return sh


def title(slide, head, sub=None, *, sub_y=1.5, h=0.95):
    text(slide, head, M, 0.5, W, h, size=32, color=INK, font=H, bold=True,
         spacing=1.05)
    if sub:
        text(slide, sub, M, sub_y, W - 0.8, 0.6, size=14.5, color=MUTE)


def card(slide, x, y, w, h, head, body, *, n=None, accent=MINT, body_size=12.5):
    panel(slide, x, y, w, h)
    hx = x + 0.32
    if n is not None:
        c = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x + 0.3),
                                   Inches(y + 0.27), Inches(0.42), Inches(0.42))
        c.fill.solid()
        c.fill.fore_color.rgb = rgb(accent)
        c.line.fill.background()
        c.shadow.inherit = False
        tf = c.text_frame
        tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        r = p.add_run()
        r.text = str(n)
        r.font.size = Pt(13)
        r.font.bold = True
        r.font.name = H
        r.font.color.rgb = rgb(BG)
        hx = x + 0.88
    text(slide, head, hx, y + 0.3, w - (hx - x) - 0.32, 0.45,
         size=14.5, color=INK, font=H, bold=True)
    text(slide, body, x + 0.32, y + 0.88, w - 0.64, h - 1.15,
         size=body_size, color=MUTE, spacing=1.18, space_after=5)


def stat(slide, x, y, w, big, label, *, color=MINT, size=36):
    text(slide, big, x, y, w, 0.66, size=size, color=color, font=H, bold=True)
    text(slide, label, x, y + 0.72, w, 0.7, size=11.5, color=MUTE, spacing=1.1)


# ------------------------------------------------------------------ 1
s = new_slide()
text(s, "P E D A L S H I E L D", M, 2.2, W, 0.4, size=15, color=MINT,
     font=H, bold=True)
text(s, "Prove the ride.\nNever the route.", M, 2.75, 9.0, 2.0,
     size=44, color=INK, font=H, bold=True, spacing=1.0)
text(s, "Verified proof that a person physically made a trip — without "
        "collecting anywhere they went.",
     M, 4.95, 9.2, 0.7, size=16, color=MUTE)
text(s, "IntelliGrip Industries   ·   Bend, Oregon   ·   pedalshield.app",
     M, 6.45, W, 0.4, size=12, color=DIM)
s.notes_slide.notes_text_frame.text = (
    "We sell verified physical-trip data to organizations legally required to "
    "measure commuting, who do it today with self-reported surveys.")

# ------------------------------------------------------------------ 2
s = new_slide()
title(s, "Nobody can prove a trip happened\nwithout surveilling the traveller.",
      "Every system that rewards or reports physical movement faces the same fork.",
      sub_y=2.05, h=1.5)
card(s, M, 2.85, 5.9, 2.8, "Trust the self-report",
     "A survey, or an app where someone taps “I biked today.”\n\n"
     "Cheap, and unverifiable. Nobody knows the real number. Incentives pay out "
     "on the honor system, and the data is weak evidence for funding or "
     "climate reporting.")
card(s, M + 6.35, 2.85, 5.88, 2.8, "Collect location data",
     "Continuous GPS from every participant.\n\n"
     "Accurate — and now you hold a database of where your employees or "
     "students go. Legal objects, participation drops, and it is a breach "
     "waiting to happen.", accent=PURP)
text(s, "Pedalshield collapses the fork: the phone does the judging, and only "
        "the verdict leaves.",
     M, 6.0, W, 0.5, size=17, color=MINT, font=H, bold=True)
s.notes_slide.notes_text_frame.text = (
    "This is the whole thesis. Not bike rewards — verification without surveillance.")

# ------------------------------------------------------------------ 3
s = new_slide()
title(s, "The demand is written into law",
      "Washington requires worksites of 100+ employees to measure drive-alone "
      "rates every two years. Oregon and California have their own mandates.")
stat(s, M, 2.5, 3.6, "1,000+",
     "WA worksites under the\nCommute Trip Reduction mandate")
stat(s, M + 4.0, 2.5, 3.6, "2027",
     "California SB 253 adds employee-commuting\nemissions to mandatory disclosure",
     color=PURP)
stat(s, M + 8.0, 2.5, 3.9, "2 min",
     "the self-reported survey that\nproduces both numbers today")
panel(s, M, 4.5, W, 1.85)
text(s, "We don’t have to create the demand, and we don’t sell one worksite at a time",
     M + 0.4, 4.72, W - 0.8, 0.4, size=15, color=INK, font=H, bold=True)
text(s, "~4,000 US worksites sit under an active commute mandate today. From 2027, "
        "every company over $1B revenue doing business in California must "
        "disclose employee-commuting emissions — a far larger set, reporting the "
        "same unverified survey. We reach both through the channels that already "
        "serve them: commute-software vendors whose weakest link is verification, "
        "and regional transportation programs that run the incentives.",
     M + 0.4, 5.2, W - 0.8, 1.05, size=12.5, spacing=1.16)
s.notes_slide.notes_text_frame.text = (
    "We are not asking anyone to believe in privacy or cycling. We offer a better "
    "number for a report they are legally required to file. First pilot target: "
    "Commute Options for Central Oregon, the region's TDM program.")

# ------------------------------------------------------------------ 4
s = new_slide()
title(s, "What we built",
      "The rider’s phone verifies the ride. A signed verdict leaves. The route never does.")
card(s, M, 2.45, 3.85, 2.65, "Sense",
     "GPS, accelerometer, gyroscope and barometer are read on-device during "
     "the ride. Nothing is uploaded.", n=1)
card(s, M + 4.05, 2.45, 3.85, 2.65, "Judge",
     "A physics engine scores the ride against what a bicycle can actually do — "
     "speed envelopes, pedaling cadence, road vibration, sensor agreement.", n=2)
card(s, M + 8.1, 2.45, 3.83, 2.65, "Sign",
     "The phone emits a distance and an integrity score, signed by a device "
     "key. That is the entire payload.", n=3)
text(s, "An open-source test fails the build if any location or sensor field "
        "could appear in that payload. The privacy promise is executable, not a "
        "policy document.",
     M, 5.45, W, 0.6, size=13.5)
s.notes_slide.notes_text_frame.text = (
    "MIT-licensed privacy contract. Anyone can clone the repo and run the test.")

# ------------------------------------------------------------------ 5
s = new_slide()
title(s, "Why this is hard — and why it is the moat",
      "Anyone can claim a trip. Making the claim expensive to fake is the entire product.")
card(s, M, 2.45, 5.9, 1.7, "Single-sensor checks fall to single-sensor fakes",
     "Replay a real GPS track while shaking the phone, and every conventional "
     "check passes.")
card(s, M, 4.3, 5.9, 2.0, "So we check that the sensors agree",
     "The gyroscope must rotate when GPS says you turned. Road vibration must "
     "rise as you accelerate. Faking one stream is easy; faking two that stay "
     "coherent in time is a different problem.", accent=PURP)
card(s, M + 6.35, 2.45, 5.88, 1.7, "Wrong answers cost more than missed ones",
     "A rejected honest rider is churn. Impossible stretches are excised and "
     "earn nothing; the rest of the ride still pays.")
card(s, M + 6.35, 4.3, 5.88, 2.0, "The engine is the licensable asset",
     "The privacy contract is open source. The scoring engine is not. Its value "
     "compounds with real ride data — the part a competitor cannot copy.")
s.notes_slide.notes_text_frame.text = (
    "Threat model is published, including current gaps. Hardware attestation in progress.")

# ------------------------------------------------------------------ 6
s = new_slide()
title(s, "It works, on mainnet, with no human in the loop",
      "Every payout is built, proved, signed and broadcast autonomously. "
      "Verify any of these yourself.")
rows = [
    ("fbf4e134…d16ed8", "First payout after the Zcash hard fork — a cross-pool migration spend"),
    ("c42a0097…5f72b8", "3.07 verified miles  ·  integrity 0.71  ·  30,000 zatoshi"),
    ("7bcda2d6…803164", "1.04 verified miles  ·  integrity 0.69  ·  17,674 zatoshi"),
]
for i, (tx, desc) in enumerate(rows):
    y = 2.45 + i * 0.92
    panel(s, M, y, W, 0.74)
    text(s, tx, M + 0.35, y, 3.3, 0.74, size=12.5, color=MINT,
         font="Courier New", bold=True, anchor=MSO_ANCHOR.MIDDLE)
    text(s, desc, M + 3.85, y, W - 4.25, 0.74, size=12.5,
         anchor=MSO_ANCHOR.MIDDLE)
stat(s, M, 5.55, 4.1, "4 min → sec",
     "payout latency, after rebuilding\nthe treasury’s chain scanner", size=30)
stat(s, M + 4.5, 5.55, 3.8, "0.62",
     "verification threshold — tuned\nagainst real rejected rides",
     color=PURP, size=30)
stat(s, M + 8.6, 5.55, 3.3, "~$0.09",
     "per mile, pegged to the EPA\nsocial cost of carbon", size=30)
s.notes_slide.notes_text_frame.text = (
    "The hard fork broke every wallet holding pre-fork funds. We restored payouts "
    "the same day with a cross-pool migration spend nobody had documented.")

# ------------------------------------------------------------------ 7
s = new_slide()
title(s, "Market — built bottoms-up from a verified anchor",
      "Washington’s 1,000+ mandated worksites sit in counties holding ~6.5M of "
      "340M Americans. Scaled by population: ~50,000 US worksites meet "
      "CTR-style criteria.")
bars = [
    ("Mandated worksites today", "4,000 × $6k", 24, MINT),
    ("Universities with TDM programs", "700 × $15k", 11, MINT),
    ("CTR-eligible US worksites", "50,000 × $6k", 300, MINT),
]
BAR_X, BAR_W, MAXV = M + 3.9, 4.0, 300.0
for i, (label, calc, val, col) in enumerate(bars):
    y = 2.75 + i * 1.05
    text(s, label, M, y - 0.04, 3.7, 0.32, size=12.5, color=INK, font=H, bold=True)
    text(s, calc, M, y + 0.3, 3.7, 0.3, size=11, color=DIM)
    panel(s, BAR_X, y + 0.02, BAR_W, 0.34, fill=EDGE, line=EDGE, radius=0.35)
    w = max(0.12, BAR_W * (val / MAXV))
    panel(s, BAR_X, y + 0.02, w, 0.34, fill=col, line=col, radius=0.35)
    text(s, f"${val}M", BAR_X + BAR_W + 0.18, y + 0.02, 1.1, 0.34,
         size=13, color=INK, font=H, bold=True, anchor=MSO_ANCHOR.MIDDLE)

panel(s, M + 9.35, 2.6, 2.88, 3.2)
text(s, "$300M", M + 9.65, 2.95, 2.3, 0.75, size=36, color=MINT, font=H, bold=True)
text(s, "beachhead TAM, per year", M + 9.65, 3.72, 2.3, 0.35, size=11.5)
text(s, "$6,000 ACV assumes 500 employees at $12/employee/year — priced as "
        "compliance software, because per-trip pricing nets only ~$1,400 per site.",
     M + 9.65, 4.2, 2.3, 1.4, size=11, spacing=1.15)
text(s, "Not sized here: from 2027 California SB 253 requires every company over "
        "$1B revenue operating in the state to disclose employee-commuting "
        "emissions. Same product, separate and larger buyer set — we have not "
        "put a number on it because we cannot yet verify the company count.",
     M, 6.2, W, 0.6, size=11.5, color=DIM, spacing=1.15)
s.notes_slide.notes_text_frame.text = (
    "Sources: WSDOT Commute Trip Reduction; Oregon DEQ ECO; SCAQMD Rule 2202.")

# ------------------------------------------------------------------ 8
s = new_slide()
title(s, "The primitive is horizontal",
      "Proof someone physically did something, without collecting where. "
      "Commute is the wedge, not the ceiling.")
card(s, M, 2.45, 3.85, 3.3, "Commute compliance",
     "~$300M/yr\n\nMandated measurement, existing budgets, and no incumbent "
     "doing verification. Our beachhead.", n=1, body_size=12)
card(s, M + 4.05, 2.45, 3.85, 3.3, "Usage-based insurance",
     "$6.9B market in 2026, ~20% of US auto policies and rising to 30–35%.\n\n"
     "Built entirely on location harvesting — and four states have introduced "
     "bills restricting it. A privacy-preserving layer at 3–5% is $200–350M/yr.",
     n=2, accent=PURP, body_size=12)
card(s, M + 8.1, 2.45, 3.83, 3.3, "Gig & delivery",
     "DoorDash lost $2.5M to a single GPS-spoofing scheme, and in 2026 banned a "
     "driver for AI-fabricated delivery evidence.\n\nPlatforms already buy "
     "location-integrity tooling.", n=3, body_size=12)
text(s, "Scale comes from licensing the engine into channels that already own "
        "the customer — one commute platform reaches thousands of enterprises; "
        "one insurance carrier is worth a thousand worksites.",
     M, 6.05, W, 0.6, size=13.5)

# ------------------------------------------------------------------ 9
s = new_slide()
title(s, "How we make money", "Two lines, and neither of them charges the rider.")
card(s, M, 2.45, 5.9, 2.55, "Compliance software",
     "A per-worksite annual contract. The deliverable is verified trip counts a "
     "program manager can put in a statutory report — replacing a "
     "self-reported survey they already know is weak.")
card(s, M + 6.35, 2.45, 5.88, 2.55, "Engine licensing",
     "The verification engine as an SDK, priced per verified trip or per seat, "
     "sold into platforms that already have the customers: commute software, "
     "telematics, delivery.", accent=PURP)
panel(s, M, 5.25, W, 1.55)
text(s, "We sell measurement, never offsets.",
     M + 0.4, 5.45, W - 0.8, 0.35, size=14, color=INK, font=H, bold=True)
text(s, "Buyers get a verified number for their own emissions inventory — not a "
        "tradable credit. The rider is already paid the carbon value of the mile, "
        "so selling a credit on the same mile would double-count it. Riders never "
        "pay, and never sell their data by default; if they opt in to contribute "
        "route data, the co-op pays them a share of what it is licensed for.",
     M + 0.4, 5.85, W - 0.8, 0.9, size=12.5, spacing=1.15)

# ----------------------------------------------------------------- 10
s = new_slide()
title(s, "Founder", "Solo. One pair of hands from the Orchard wallet to the App Store build.")
panel(s, M, 2.4, 5.9, 3.9)
text(s, "Samuel B. Newman", M + 0.4, 2.7, 5.1, 0.45,
     size=22, color=INK, font=H, bold=True)
text(s, "Founder · IntelliGrip Industries · Bend, Oregon",
     M + 0.4, 3.2, 5.1, 0.35, size=13, color=MINT)
text(s, "Before this I built eBikes Tubes Direct, an e-commerce tube and tire "
        "company I scaled to $300K — selling to cyclists, one order at a time.\n"
        "\n"
        "A verified bicycle mile is worth nothing today. Neither is the "
        "pollution it avoided. Not because the value isn’t real, but because "
        "nobody can prove the mile happened.\n"
        "\n"
        "I am building that proof.",
     M + 0.4, 3.72, 5.1, 2.45, size=13, spacing=1.22)

stat(s, M + 6.7, 2.4, 2.5, "10 weeks",
     "from first commit to a live\niOS app paying on mainnet", size=28)
stat(s, M + 9.6, 2.4, 2.6, "100", "commits, one author", size=28, color=PURP)
stat(s, M + 6.7, 4.15, 2.5, "~40k",
     "lines of Rust and\nTypeScript, shipped", size=28, color=PURP)
stat(s, M + 9.6, 4.15, 2.6, "$300K", "revenue at my previous\ncompany, also solo", size=28)

panel(s, M, 6.5, W, 0.75)
text(s, "Built solo: a hand-rolled Zcash Orchard wallet, an on-device anti-cheat "
        "engine, an autonomous payout treasury, and an iOS app — then kept it "
        "paying through a consensus hard fork the same day it broke.",
     M + 0.4, 6.68, W - 0.8, 0.45, size=12.5)
s.notes_slide.notes_text_frame.text = (
    "The hard fork is the story worth telling live: legacy-pool funds could no "
    "longer be sent to anyone. I built a cross-pool migration spend that carries "
    "the migration inside each payout, so my existing scanner kept working.")

# ----------------------------------------------------------------- 11
s = new_slide()
title(s, "Where we actually are", "Stated plainly, because the gap is the point.")
card(s, M, 2.45, 5.9, 2.6, "Built and live",
     "iOS app shipping on TestFlight\n"
     "Autonomous shielded payouts on mainnet, no operator\n"
     "An anti-cheat engine that survived a consensus hard fork\n"
     "Published threat model, MIT privacy contract")
card(s, M + 6.35, 2.45, 5.88, 2.6, "Not yet",
     "One rider — the founder\n"
     "Zero revenue, zero customers\n"
     "Thresholds tuned against synthetic attacks and a handful of real rides\n"
     "No pilot signed", accent=PURP)
panel(s, M, 5.35, W, 1.4, line=MINT, line_w=1.5)
text(s, "The ask: ten riders and one signed pilot in Bend.",
     M + 0.4, 5.57, W - 0.8, 0.42, size=17, color=MINT, font=H, bold=True)
text(s, "That pilot answers the only question no amount of engineering can: how "
        "far apart are self-reported and verified trip counts — and will an "
        "organization pay to know.",
     M + 0.4, 6.05, W - 0.8, 0.6, size=13)
s.notes_slide.notes_text_frame.text = (
    "Pilot targets: Commute Options for Central Oregon, OSU-Cascades Mobility "
    "Lab, City of Bend.")

out = sys.argv[1] if len(sys.argv) > 1 else "Pedalshield_Deck.pptx"
prs.save(out)
print("wrote", out, "-", len(prs.slides.__iter__.__self__._sldIdLst), "slides")
